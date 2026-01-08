from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RgbColor(26, 26, 46)
BLUE = RgbColor(66, 133, 244)
GREEN = RgbColor(52, 168, 83)
PURPLE = RgbColor(155, 89, 182)
RED = RgbColor(234, 67, 53)
WHITE = RgbColor(255, 255, 255)
GRAY = RgbColor(136, 146, 176)

def add_background(slide, r, g, b):
    """Add solid color background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(r, g, b)
    background.line.fill.background()
    # Move to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_slide_number(slide, num, total=7):
    """Add slide number to bottom right"""
    txBox = slide.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_bullet_item(slide, left, top, width, emoji, title, description, accent_color):
    """Add a styled bullet point"""
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(40, 40, 60)
    shape.line.fill.background()
    
    # Accent line on left
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(0.7)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    # Emoji
    emoji_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), Inches(0.5), Inches(0.5))
    emoji_tf = emoji_box.text_frame
    emoji_p = emoji_tf.paragraphs[0]
    emoji_p.text = emoji
    emoji_p.font.size = Pt(24)
    
    # Text
    text_box = slide.shapes.add_textbox(left + Inches(0.7), top + Inches(0.15), width - Inches(0.9), Inches(0.5))
    text_tf = text_box.text_frame
    text_tf.word_wrap = True
    p = text_tf.paragraphs[0]
    
    run1 = p.add_run()
    run1.text = title + " "
    run1.font.bold = True
    run1.font.size = Pt(18)
    run1.font.color.rgb = accent_color
    
    run2 = p.add_run()
    run2.text = "— " + description
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE

# ========== SLIDE 1: Title ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide1, 15, 12, 41)

# Cloud emoji
cloud = slide1.shapes.add_textbox(Inches(5.8), Inches(1), Inches(2), Inches(1.2))
cloud_tf = cloud.text_frame
cloud_p = cloud_tf.paragraphs[0]
cloud_p.text = "☁️"
cloud_p.font.size = Pt(80)
cloud_p.alignment = PP_ALIGN.CENTER

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "Використання хмарних сервісів"
title_p.font.size = Pt(48)
title_p.font.bold = True
title_p.font.color.rgb = BLUE
title_p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(0.6))
sub_tf = sub_box.text_frame
sub_p = sub_tf.paragraphs[0]
sub_p.text = "можливості, переваги та ризики"
sub_p.font.size = Pt(28)
sub_p.font.color.rgb = GRAY
sub_p.alignment = PP_ALIGN.CENTER

# Author info
author_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
author_tf = author_box.text_frame
author_p = author_tf.paragraphs[0]
author_p.text = "👤 Автор: Студент"
author_p.font.size = Pt(18)
author_p.font.color.rgb = GRAY
author_p.alignment = PP_ALIGN.CENTER

date_p = author_tf.add_paragraph()
date_p.text = "📅 Січень 2026"
date_p.font.size = Pt(18)
date_p.font.color.rgb = GRAY
date_p.alignment = PP_ALIGN.CENTER

add_slide_number(slide1, 1)

# ========== SLIDE 2: What are cloud services ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, 26, 26, 46)

# Title
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
title2_tf = title2.text_frame
title2_p = title2_tf.paragraphs[0]
title2_p.text = "☁️  Що таке хмарні сервіси?"
title2_p.font.size = Pt(40)
title2_p.font.bold = True
title2_p.font.color.rgb = BLUE

# Bullet points
bullets2 = [
    ("🌐", "Доступ через інтернет", "працюйте з файлами з будь-якого пристрою"),
    ("💾", "Віддалене зберігання", "дані на серверах, не на вашому комп'ютері"),
    ("⚡", "Обчислення за запитом", "потужні сервери виконують складні задачі"),
    ("🔄", "Синхронізація", "автоматичне оновлення на всіх пристроях"),
]

for i, (emoji, title, desc) in enumerate(bullets2):
    add_bullet_item(slide2, Inches(0.8), Inches(1.8 + i * 1.1), Inches(11.5), emoji, title, desc, BLUE)

add_slide_number(slide2, 2)

# ========== SLIDE 3: Advantages ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, 13, 27, 42)

title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
title3_tf = title3.text_frame
title3_p = title3_tf.paragraphs[0]
title3_p.text = "✅  Переваги хмарних сервісів"
title3_p.font.size = Pt(40)
title3_p.font.bold = True
title3_p.font.color.rgb = GREEN

bullets3 = [
    ("🕐", "Доступність 24/7", "працюйте будь-коли та будь-де"),
    ("👥", "Спільна робота", "редагуйте документи разом у реальному часі"),
    ("🛡️", "Безпека даних", "автоматичне резервне копіювання"),
    ("📈", "Масштабованість", "легко збільшити обсяг пам'яті"),
    ("💰", "Економія", "не потрібно купувати сервери"),
]

for i, (emoji, title, desc) in enumerate(bullets3):
    add_bullet_item(slide3, Inches(0.8), Inches(1.5 + i * 1.0), Inches(11.5), emoji, title, desc, GREEN)

add_slide_number(slide3, 3)

# ========== SLIDE 4: Education ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, 26, 26, 58)

title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
title4_tf = title4.text_frame
title4_p = title4_tf.paragraphs[0]
title4_p.text = "🎓  Хмарні сервіси в освіті"
title4_p.font.size = Pt(40)
title4_p.font.bold = True
title4_p.font.color.rgb = PURPLE

bullets4 = [
    ("📝", "Google Docs", "спільні документи для групових проєктів"),
    ("📋", "Google Forms", "онлайн-тестування та опитування"),
    ("🏫", "Google Classroom", "віртуальні класи та завдання"),
    ("💬", "Microsoft Teams", "відеоконференції та командна робота"),
]

for i, (emoji, title, desc) in enumerate(bullets4):
    add_bullet_item(slide4, Inches(0.8), Inches(1.8 + i * 1.1), Inches(11.5), emoji, title, desc, PURPLE)

add_slide_number(slide4, 4)

# ========== SLIDE 5: Risks ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, 35, 26, 26)

title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
title5_tf = title5.text_frame
title5_p = title5_tf.paragraphs[0]
title5_p.text = "⚠️  Ризики та загрози"
title5_p.font.size = Pt(40)
title5_p.font.bold = True
title5_p.font.color.rgb = RED

bullets5 = [
    ("📡", "Залежність від інтернету", "без з'єднання немає доступу"),
    ("🔓", "Конфіденційність", "ризик витоку особистих даних"),
    ("💥", "Технічні збої", "можлива втрата даних при аваріях"),
    ("🏢", "Залежність від провайдера", "закриття сервісу = втрата доступу"),
]

for i, (emoji, title, desc) in enumerate(bullets5):
    add_bullet_item(slide5, Inches(0.8), Inches(1.8 + i * 1.1), Inches(11.5), emoji, title, desc, RED)

add_slide_number(slide5, 5)

# ========== SLIDE 6: Chart ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, 26, 26, 46)

title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
title6_tf = title6.text_frame
title6_p = title6_tf.paragraphs[0]
title6_p.text = "📊  Динаміка використання 2010–2025"
title6_p.font.size = Pt(40)
title6_p.font.bold = True
title6_p.font.color.rgb = BLUE

# Chart container
chart_bg = slide6.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5)
)
chart_bg.fill.solid()
chart_bg.fill.fore_color.rgb = RgbColor(30, 30, 50)
chart_bg.line.fill.background()

# Bar chart data
chart_data = [
    ("2010", 45, 15),
    ("2012", 105, 35),
    ("2014", 285, 70),
    ("2016", 695, 120),
    ("2018", 1450, 180),
    ("2020", 3200, 280),
    ("2022", 5280, 350),
    ("2024", 7100, 420),
    ("2025", 8200, 480),
]

bar_width = Inches(0.8)
chart_left = Inches(1.5)
chart_bottom = Inches(6.3)
max_height = Inches(4.2)

for i, (year, value, height_px) in enumerate(chart_data):
    bar_height = Inches(height_px / 120)
    bar_left = chart_left + Inches(i * 1.2)
    bar_top = chart_bottom - bar_height
    
    # Bar
    bar = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE if i < 5 else GREEN
    bar.line.fill.background()
    
    # Value label
    val_box = slide6.shapes.add_textbox(bar_left - Inches(0.1), bar_top - Inches(0.35), Inches(1), Inches(0.3))
    val_tf = val_box.text_frame
    val_p = val_tf.paragraphs[0]
    if value >= 1000:
        val_p.text = f"{value/1000:.1f}B"
    else:
        val_p.text = f"{value}M"
    val_p.font.size = Pt(11)
    val_p.font.color.rgb = BLUE if i < 5 else GREEN
    val_p.alignment = PP_ALIGN.CENTER
    
    # Year label
    year_box = slide6.shapes.add_textbox(bar_left, chart_bottom + Inches(0.05), bar_width, Inches(0.3))
    year_tf = year_box.text_frame
    year_p = year_tf.paragraphs[0]
    year_p.text = year
    year_p.font.size = Pt(12)
    year_p.font.color.rgb = GRAY
    year_p.alignment = PP_ALIGN.CENTER

# Bottom note
note_box = slide6.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
note_tf = note_box.text_frame
note_p = note_tf.paragraphs[0]
note_p.text = "📈 Зростання з 45 млн до 8.2 млрд користувачів за 15 років"
note_p.font.size = Pt(16)
note_p.font.color.rgb = GRAY
note_p.alignment = PP_ALIGN.CENTER

add_slide_number(slide6, 6)

# ========== SLIDE 7: Links ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, 15, 12, 41)

title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
title7_tf = title7.text_frame
title7_p = title7_tf.paragraphs[0]
title7_p.text = "🔗  Дізнайтеся більше"
title7_p.font.size = Pt(40)
title7_p.font.bold = True
title7_p.font.color.rgb = BLUE

# Link cards
cards = [
    ("🎬", "Відео-урок", "Що таке хмарні технології?\nПростими словами", "youtube.com"),
    ("📚", "Документація AWS", "Офіційний посібник з\nхмарних обчислень", "aws.amazon.com"),
    ("🎓", "Курс Coursera", "Introduction to\nCloud Computing", "coursera.org"),
]

card_width = Inches(3.8)
card_height = Inches(4.2)
card_start = Inches(0.8)
card_gap = Inches(0.3)

for i, (emoji, title, desc, url) in enumerate(cards):
    card_left = card_start + i * (card_width + card_gap)
    
    # Card background
    card = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, card_left, Inches(1.8), card_width, card_height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RgbColor(35, 35, 55)
    card.line.color.rgb = RgbColor(60, 60, 80)
    
    # Emoji
    emoji_box = slide7.shapes.add_textbox(card_left, Inches(2.1), card_width, Inches(0.8))
    emoji_tf = emoji_box.text_frame
    emoji_p = emoji_tf.paragraphs[0]
    emoji_p.text = emoji
    emoji_p.font.size = Pt(50)
    emoji_p.alignment = PP_ALIGN.CENTER
    
    # Title
    card_title = slide7.shapes.add_textbox(card_left + Inches(0.2), Inches(3.0), card_width - Inches(0.4), Inches(0.5))
    card_title_tf = card_title.text_frame
    card_title_p = card_title_tf.paragraphs[0]
    card_title_p.text = title
    card_title_p.font.size = Pt(20)
    card_title_p.font.bold = True
    card_title_p.font.color.rgb = BLUE
    card_title_p.alignment = PP_ALIGN.CENTER
    
    # Description
    card_desc = slide7.shapes.add_textbox(card_left + Inches(0.2), Inches(3.6), card_width - Inches(0.4), Inches(1))
    card_desc_tf = card_desc.text_frame
    card_desc_tf.word_wrap = True
    card_desc_p = card_desc_tf.paragraphs[0]
    card_desc_p.text = desc
    card_desc_p.font.size = Pt(14)
    card_desc_p.font.color.rgb = GRAY
    card_desc_p.alignment = PP_ALIGN.CENTER
    
    # URL button
    btn = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, card_left + Inches(0.5), Inches(5.1), card_width - Inches(1), Inches(0.5)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = BLUE
    btn.line.fill.background()
    
    btn_text = slide7.shapes.add_textbox(card_left + Inches(0.5), Inches(5.15), card_width - Inches(1), Inches(0.45))
    btn_tf = btn_text.text_frame
    btn_p = btn_tf.paragraphs[0]
    btn_p.text = url + " →"
    btn_p.font.size = Pt(12)
    btn_p.font.bold = True
    btn_p.font.color.rgb = WHITE
    btn_p.alignment = PP_ALIGN.CENTER

add_slide_number(slide7, 7)

# Save presentation
output_path = "/Users/maksymyanisiv/test for cursor/cloud_services_presentation.pptx"
prs.save(output_path)
print(f"✅ Презентацію збережено: {output_path}")

